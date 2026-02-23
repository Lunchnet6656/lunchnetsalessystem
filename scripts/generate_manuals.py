"""
シフト管理システム 操作マニュアル PPTX 生成スクリプト
使い方: python scripts/generate_manuals.py
出力:   employee_manual.pptx (パート従業員向け)
        admin_manual.pptx    (管理者向け)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ---------------------------------------------------------------------------
# デザイン定数
# ---------------------------------------------------------------------------
TITLE_COLOR   = RGBColor(0x1E, 0x40, 0xAF)  # 青
SUBTITLE_COLOR= RGBColor(0x37, 0x4F, 0x6D)  # 濃紺グレー
ACCENT_GREEN  = RGBColor(0x16, 0xA3, 0x4A)  # 緑
ACCENT_RED    = RGBColor(0xDC, 0x26, 0x26)  # 赤
ACCENT_YELLOW = RGBColor(0xCA, 0x8A, 0x04)  # 黄
ACCENT_ORANGE = RGBColor(0xEA, 0x58, 0x0C)  # オレンジ
BG_BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)  # 薄青
BG_GRAY       = RGBColor(0xF3, 0xF4, 0xF6)  # 薄グレー
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK         = RGBColor(0x11, 0x18, 0x27)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# 共通ヘルパー
# ---------------------------------------------------------------------------

def _set_font(run, size=None, bold=False, color=None, name="Meiryo UI"):
    run.font.name = name
    if size:
        run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _add_textbox(slide, text, left, top, width, height,
                 font_size=16, bold=False, color=BLACK,
                 align=PP_ALIGN.LEFT, word_wrap=True, font_name="Meiryo UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, font_size, bold, color, font_name)
    return txBox


def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _header_bar(slide, title_text, color=TITLE_COLOR):
    """スライド上部の色帯 + タイトルテキスト"""
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), color)
    _add_textbox(slide, title_text,
                 Inches(0.4), Inches(0.12),
                 Inches(12.5), Inches(0.9),
                 font_size=28, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT)


def _footer(slide, text="シフト管理システム 操作マニュアル"):
    _add_rect(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4),
              RGBColor(0xE5, 0xE7, 0xEB))
    _add_textbox(slide, text,
                 Inches(0.3), Inches(7.1),
                 Inches(12.7), Inches(0.4),
                 font_size=10, color=RGBColor(0x6B, 0x72, 0x80),
                 align=PP_ALIGN.LEFT)


def add_title_slide(prs, title, subtitle, color=TITLE_COLOR):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 背景グラデーション風 (上部矩形 + 下部白)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_GRAY)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(3.5), color)

    _add_textbox(slide, title,
                 Inches(1), Inches(0.7),
                 Inches(11.3), Inches(2.2),
                 font_size=36, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
    _add_textbox(slide, subtitle,
                 Inches(1), Inches(3.8),
                 Inches(11.3), Inches(1.0),
                 font_size=22, bold=False, color=SUBTITLE_COLOR,
                 align=PP_ALIGN.CENTER)
    _add_textbox(slide, "2026年2月改訂",
                 Inches(1), Inches(5.0),
                 Inches(11.3), Inches(0.5),
                 font_size=14, color=RGBColor(0x6B, 0x72, 0x80),
                 align=PP_ALIGN.CENTER)
    return slide


def add_toc_slide(prs, items, color=TITLE_COLOR):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_GRAY)
    _header_bar(slide, "目次", color)

    cols = 2
    col_w = Inches(5.8)
    start_x = [Inches(0.5), Inches(7.0)]
    y = Inches(1.3)
    line_h = Inches(0.42)

    for i, item in enumerate(items):
        col = i % cols
        row = i // cols
        x = start_x[col]
        top = y + row * line_h
        num = f"{i+1:02d}. "
        _add_textbox(slide, num + item, x, top, col_w, line_h,
                     font_size=15, color=BLACK)

    _footer(slide)
    return slide


def add_content_slide(prs, title, bullets, color=TITLE_COLOR, numbered=False):
    """タイトル + 箇条書きスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_GRAY)
    _header_bar(slide, title, color)

    y = Inches(1.25)
    for item in bullets:
        # ネスト判定
        if isinstance(item, tuple):
            indent, text = item
        else:
            indent, text = 0, item

        indent_x = Inches(0.5 + indent * 0.4)
        bullet_char = "●" if indent == 0 else "◦"
        if numbered and indent == 0:
            bullet_char = ""

        _add_textbox(slide, f"  {bullet_char}  {text}" if bullet_char else f"  {text}",
                     indent_x, y,
                     Inches(12.3 - indent * 0.4), Inches(0.5),
                     font_size=16 if indent == 0 else 14,
                     color=BLACK if indent > 0 else BLACK)
        y += Inches(0.45)

    _footer(slide)
    return slide


def add_steps_slide(prs, title, steps, color=TITLE_COLOR):
    """番号付き手順スライド（丸囲み数字 + 説明）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_GRAY)
    _header_bar(slide, title, color)

    y = Inches(1.3)
    step_h = Inches(0.75)
    for i, (step_title, step_desc) in enumerate(steps, 1):
        # 番号丸
        circle = _add_rect(slide, Inches(0.4), y + Inches(0.05),
                            Inches(0.55), Inches(0.55), color)
        _add_textbox(slide, str(i),
                     Inches(0.4), y + Inches(0.05),
                     Inches(0.55), Inches(0.55),
                     font_size=16, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        # 矢印線
        if i < len(steps):
            _add_rect(slide,
                      Inches(0.62), y + Inches(0.6),
                      Inches(0.12), Inches(0.25),
                      RGBColor(0x9C, 0xA3, 0xAF))

        # テキスト
        _add_textbox(slide, step_title,
                     Inches(1.1), y,
                     Inches(11.5), Inches(0.35),
                     font_size=16, bold=True, color=color)
        _add_textbox(slide, step_desc,
                     Inches(1.1), y + Inches(0.35),
                     Inches(11.5), Inches(0.35),
                     font_size=13, color=RGBColor(0x37, 0x4F, 0x6D))
        y += step_h

    _footer(slide)
    return slide


def add_two_column_slide(prs, title, left_title, left_items,
                         right_title, right_items, color=TITLE_COLOR):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_GRAY)
    _header_bar(slide, title, color)

    mid = Inches(6.7)
    # 左列背景
    _add_rect(slide, Inches(0.3), Inches(1.2),
              Inches(6.0), Inches(5.7), WHITE)
    # 右列背景
    _add_rect(slide, Inches(7.0), Inches(1.2),
              Inches(6.0), Inches(5.7), WHITE)

    # 左タイトル
    _add_rect(slide, Inches(0.3), Inches(1.2),
              Inches(6.0), Inches(0.45), color)
    _add_textbox(slide, left_title,
                 Inches(0.5), Inches(1.22),
                 Inches(5.6), Inches(0.4),
                 font_size=15, bold=True, color=WHITE)
    # 右タイトル
    _add_rect(slide, Inches(7.0), Inches(1.2),
              Inches(6.0), Inches(0.45), color)
    _add_textbox(slide, right_title,
                 Inches(7.2), Inches(1.22),
                 Inches(5.6), Inches(0.4),
                 font_size=15, bold=True, color=WHITE)

    y_l = Inches(1.75)
    for item in left_items:
        _add_textbox(slide, f"● {item}", Inches(0.5), y_l,
                     Inches(5.6), Inches(0.45), font_size=14, color=BLACK)
        y_l += Inches(0.42)

    y_r = Inches(1.75)
    for item in right_items:
        _add_textbox(slide, f"● {item}", Inches(7.2), y_r,
                     Inches(5.6), Inches(0.45), font_size=14, color=BLACK)
        y_r += Inches(0.42)

    _footer(slide)
    return slide


def add_heatmap_slide(prs, color=TITLE_COLOR):
    """ヒートマップ4色凡例スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_GRAY)
    _header_bar(slide, "ヒートマップの見方", color)

    entries = [
        (ACCENT_GREEN,  "緑 (OK)",     "出勤者数が目標を達成しています。そのまま提出できます。"),
        (ACCENT_YELLOW, "黄 (注意)",   "やや人数が不足しています。可能であれば出勤に変更を検討してください。"),
        (ACCENT_RED,    "赤 (危険)",   "深刻な人員不足です。可能な限り出勤への変更に協力してください。"),
        (ACCENT_ORANGE, "橙 (運転不足)", "ドライバー要員が不足しています。ドライバー資格保有者は出勤を優先してください。"),
    ]

    x_box = Inches(0.5)
    y = Inches(1.4)
    box_h = Inches(1.2)
    box_w = Inches(12.3)
    gap = Inches(0.1)

    for fill, label, desc in entries:
        _add_rect(slide, x_box, y, Inches(0.4), box_h, fill)
        _add_rect(slide, x_box + Inches(0.4), y, box_w - Inches(0.4), box_h,
                  WHITE, line_color=fill)
        _add_textbox(slide, label,
                     x_box + Inches(0.55), y + Inches(0.1),
                     Inches(3.5), Inches(0.45),
                     font_size=17, bold=True, color=fill)
        _add_textbox(slide, desc,
                     x_box + Inches(0.55), y + Inches(0.55),
                     Inches(11.5), Inches(0.55),
                     font_size=14, color=BLACK)
        y += box_h + gap

    _footer(slide)
    return slide


def add_holiday_reason_slide(prs, color=TITLE_COLOR):
    """休みの理由区分スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_GRAY)
    _header_bar(slide, "休みの理由区分", color)

    reasons = [
        ("私用",           TITLE_COLOR,
         "個人的な用事・予定による休み。理由を問われない最も一般的な区分です。"),
        ("収入調整",       ACCENT_GREEN,
         "扶養内勤務などで収入を一定額以下に抑えるための休みです。"),
        ("代理休み",       ACCENT_ORANGE,
         "他の従業員の代わりに休む場合に選択します。代替者の指定が必要です。"),
        ("冠婚葬祭・療養入院", ACCENT_RED,
         "結婚式・葬儀・入院・療養など、特別な事情による休みです。"),
    ]

    y = Inches(1.3)
    for label, lcolor, desc in reasons:
        _add_rect(slide, Inches(0.4), y, Inches(2.8), Inches(0.9), lcolor)
        _add_textbox(slide, label,
                     Inches(0.4), y + Inches(0.15),
                     Inches(2.8), Inches(0.6),
                     font_size=16, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        _add_rect(slide, Inches(3.3), y, Inches(9.7), Inches(0.9),
                  WHITE, line_color=lcolor)
        _add_textbox(slide, desc,
                     Inches(3.5), y + Inches(0.15),
                     Inches(9.3), Inches(0.6),
                     font_size=14, color=BLACK)
        y += Inches(1.1)

    _footer(slide)
    return slide


def add_status_badge_slide(prs, color=TITLE_COLOR):
    """ステータスバッジ説明スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_GRAY)
    _header_bar(slide, "期間選択画面 ─ ステータスバッジの見方", color)

    statuses = [
        ("募集中",   ACCENT_GREEN,  "現在シフト希望を受け付けています。提出期限内に提出してください。"),
        ("未提出",   ACCENT_YELLOW, "提出期限が過ぎていますが、まだ提出していません。速やかに提出を。"),
        ("再提出",   ACCENT_ORANGE, "管理者から返却されました。コメントを確認して再提出してください。"),
        ("提出済み", TITLE_COLOR,   "提出が完了しています。承認・確定を待ちましょう。"),
        ("確定済み", ACCENT_GREEN,  "シフトが確定されました。確定スケジュール画面で確認できます。"),
    ]

    y = Inches(1.35)
    for label, lcolor, desc in statuses:
        _add_rect(slide, Inches(0.4), y + Inches(0.05),
                  Inches(1.8), Inches(0.5), lcolor)
        _add_textbox(slide, label,
                     Inches(0.4), y + Inches(0.05),
                     Inches(1.8), Inches(0.5),
                     font_size=14, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        _add_textbox(slide, desc,
                     Inches(2.4), y + Inches(0.05),
                     Inches(10.5), Inches(0.5),
                     font_size=14, color=BLACK)
        y += Inches(0.72)

    _footer(slide)
    return slide


def add_period_status_slide(prs, color=TITLE_COLOR):
    """管理者向け: 期間ステータス遷移スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_GRAY)
    _header_bar(slide, "期間のステータス遷移", color)

    statuses = [
        ("OPEN",      ACCENT_GREEN,  "募集中: 従業員がシフト希望を提出できる状態"),
        ("REVIEW",    ACCENT_YELLOW, "審査中: 提出締切後、管理者が内容を確認・承認する状態"),
        ("FIXED",     ACCENT_ORANGE, "確定済み: 管理者がシフトを確定。公開前の最終確認状態"),
        ("PUBLISHED", TITLE_COLOR,   "公開済み: 確定シフトを全従業員が閲覧可能な状態"),
    ]

    # 矢印フロー
    y_center = Inches(2.5)
    box_w = Inches(2.4)
    box_h = Inches(1.0)
    arrow_w = Inches(0.6)
    total_w = len(statuses) * box_w + (len(statuses) - 1) * arrow_w
    x = (SLIDE_W - total_w) / 2

    for i, (label, lcolor, desc) in enumerate(statuses):
        _add_rect(slide, x, y_center - box_h/2, box_w, box_h, lcolor)
        _add_textbox(slide, label,
                     x, y_center - box_h/2,
                     box_w, box_h,
                     font_size=18, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        if i < len(statuses) - 1:
            # 矢印
            _add_rect(slide,
                      x + box_w, y_center - Inches(0.12),
                      arrow_w, Inches(0.24),
                      RGBColor(0x9C, 0xA3, 0xAF))
            _add_textbox(slide, "▶",
                         x + box_w, y_center - Inches(0.2),
                         arrow_w, Inches(0.4),
                         font_size=14, color=RGBColor(0x6B, 0x72, 0x80),
                         align=PP_ALIGN.CENTER)
        x += box_w + arrow_w

    # 説明テキスト
    y = Inches(3.7)
    x = (SLIDE_W - total_w) / 2
    for i, (label, lcolor, desc) in enumerate(statuses):
        _add_textbox(slide, desc,
                     x, y, box_w + arrow_w, Inches(1.2),
                     font_size=12, color=BLACK,
                     align=PP_ALIGN.CENTER, word_wrap=True)
        x += box_w + arrow_w

    _footer(slide)
    return slide


# ---------------------------------------------------------------------------
# 従業員向けマニュアル
# ---------------------------------------------------------------------------

def build_employee_manual(output_path="employee_manual.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    C = TITLE_COLOR  # アクセントカラー

    # 1. タイトル
    add_title_slide(prs,
                    "パート従業員向け\nシフト管理システム 操作マニュアル",
                    "～ シフト希望の提出から確定スケジュール確認まで ～")

    # 2. 目次
    toc = [
        "システム概要",
        "ログインとホーム画面",
        "シフト提出の全体フロー",
        "期間選択画面",
        "シフト希望提出画面①（画面構成）",
        "シフト希望提出画面②（日程の入力）",
        "ヒートマップの見方",
        "休みの理由区分",
        "代理休みの手順",
        "備考・コメントの入力",
        "提出ボタンと注意事項",
        "提出状況確認画面",
        "確定シフトの確認",
        "LINE連携の設定方法",
        "通知設定",
        "ルールと規則の概要",
        "よくある質問 (Q&A)",
    ]
    add_toc_slide(prs, toc)

    # 3. システム概要
    add_content_slide(prs, "システム概要", [
        "このシステムでできること:",
        (1, "2週間単位でシフト希望をオンラインで提出できます"),
        (1, "ヒートマップで職場全体の人員状況をリアルタイム確認できます"),
        (1, "管理者からの通知（LINE / メール）を受け取れます"),
        (1, "確定後のシフトスケジュールをいつでも確認できます"),
        "対象者:",
        (1, "パート・アルバイト従業員（全員）"),
        "動作環境:",
        (1, "スマートフォン・タブレット・PC のブラウザ（Chrome 推奨）"),
        (1, "社内 URL または管理者から通知された URL でアクセス"),
    ])

    # 4. ログインとホーム画面
    add_steps_slide(prs, "ログインとホーム画面", [
        ("ブラウザで URL にアクセス",
         "管理者から通知された URL を入力してください（例: https://xxxx.com/login/）"),
        ("ユーザー名 と パスワード を入力",
         "初回ログイン時のパスワードは管理者に確認してください"),
        ("「ログイン」ボタンをクリック",
         "ログインするとホーム画面（シフト提出メニュー）が表示されます"),
        ("ホーム画面を確認",
         "「シフト提出」「確定スケジュール」「通知設定」などのメニューが表示されます"),
    ])

    # 5. シフト提出の全体フロー
    add_steps_slide(prs, "シフト提出の全体フロー", [
        ("期間選択", "募集中の期間カードを選択します"),
        ("日程入力", "各日付の出勤 / 休みを入力し、理由や備考を記入します"),
        ("「提出する」ボタンをクリック", "提出完了後、ステータスが「提出済み」になります"),
        ("管理者による審査", "管理者が内容を確認し、承認または返却します"),
        ("確定スケジュール確認", "PUBLISHED 後、確定スケジュール画面で自分の割当を確認します"),
    ])

    # 6. 期間選択画面
    add_status_badge_slide(prs)

    # 7. シフト希望提出画面①（画面構成）
    add_content_slide(prs, "シフト希望提出画面①（画面構成）", [
        "画面上部に表示される情報:",
        (1, "対象期間: シフト希望を提出する2週間の開始日〜終了日"),
        (1, "提出締切: この日時までに提出しないと変更が困難になります"),
        (1, "勤務形態: フルタイム / パート など自分の区分が表示されます"),
        "カレンダー形式で日付ごとに入力欄が表示されます",
        "画面下部にヒートマップと職場全体の人員状況が表示されます",
        "重要: 締切日時を必ず確認してから入力を開始しましょう",
    ])

    # 8. シフト希望提出画面②（日程の入力）
    add_content_slide(prs, "シフト希望提出画面②（日程の入力）", [
        "各日付の入力方法:",
        (1, "「出勤」「休み」のいずれかをタップ / クリックして選択します"),
        (1, "「休み」を選択した場合、理由区分（私用 / 収入調整 など）を選びます"),
        (1, "日付ごとにコメントを入力できます（任意）"),
        "入力のコツ:",
        (1, "出勤できる日はできるだけ「出勤」を選択してください"),
        (1, "ヒートマップが赤い日は人員不足です。出勤への変更を検討してください"),
        (1, "入力途中でも自動保存されます（「下書き保存」ボタン不要）"),
        "注意: 会社指定の休日（公休）は自動的にグレーアウト表示されます",
    ])

    # 9. ヒートマップの見方
    add_heatmap_slide(prs)

    # 10. 休みの理由区分
    add_holiday_reason_slide(prs)

    # 11. 代理休みの手順
    add_steps_slide(prs, "代理休みの手順", [
        ("休みの理由で「代理休み」を選択",
         "日付ごとの理由区分から「代理休み」を選んでください"),
        ("代替者（代わりに出勤する人）を選択",
         "リストから代替者の名前を選択します。事前に当該従業員と合意してください"),
        ("備考欄に状況を記入（推奨）",
         "「〇〇さんと交代のため」などコメントを添えると管理者が確認しやすくなります"),
        ("提出する",
         "代替者にも通知が送られます。双方のシフトが連動して管理されます"),
    ])

    # 12. 備考・コメントの入力
    add_content_slide(prs, "備考・コメントの入力", [
        "2種類のコメント欄があります:",
        (1, "【日付別コメント】: 各日付の入力欄の下にあるテキスト欄"),
        (1, "    例: 「午前のみ可」「遅刻の可能性あり」など"),
        (1, "【全体備考】: ページ下部にある全体向けのコメント欄"),
        (1, "    例: 「今月は研修のため木曜は終日不可です」など"),
        "コメントは管理者のみ閲覧します。他の従業員には表示されません",
        "日本語・英語どちらでも入力可能です",
        "入力は任意ですが、管理者への伝達事項がある場合は積極的に活用してください",
    ])

    # 13. 提出ボタンと注意事項
    add_content_slide(prs, "提出ボタンと注意事項", [
        "「提出する」ボタンについて:",
        (1, "全日付の入力が完了したら、ページ下部の「提出する」ボタンをクリックします"),
        (1, "提出完了後、確認メッセージが表示されます"),
        (1, "ステータスが「提出済み」に変わったことを確認してください"),
        "締切後の変更について:",
        (1, "提出締切を過ぎると「提出する」ボタンが非表示になります"),
        (1, "締切後の修正は管理者への連絡が必要です"),
        (1, "管理者が「返却」した場合のみ、再提出が可能になります"),
        "注意: 提出前にすべての日付を入力したか必ず確認してください！",
    ])

    # 14. 提出状況確認画面
    add_content_slide(prs, "提出状況確認画面", [
        "メニューの「提出状況」から確認できます",
        "表示されるステータス:",
        (1, "提出済み: 管理者の審査を待っている状態"),
        (1, "承認済み: 管理者が内容を承認しました"),
        (1, "返却: 修正が必要です。管理者メモを確認してください"),
        "管理者メモの確認方法:",
        (1, "「返却」ステータスの場合、行をクリックすると管理者からのメモが表示されます"),
        (1, "メモの内容に従って修正し、再提出してください"),
    ])

    # 15. 確定シフトの確認
    add_content_slide(prs, "確定シフトの確認", [
        "確定スケジュール画面の開き方:",
        (1, "ホーム画面の「確定スケジュール」をクリックします"),
        (1, "確認したい期間を選択します"),
        "画面の見方:",
        (1, "自分の割当は【青太字】で強調表示されます"),
        (1, "他の従業員の割当はグレー表示されます"),
        (1, "会社の公休日はグレー背景で表示されます"),
        "注意: 確定後の変更は管理者への申請が必要です。直接変更はできません",
    ])

    # 16. LINE連携の設定方法
    add_steps_slide(prs, "LINE連携の設定方法", [
        ("通知設定画面を開く",
         "ホーム画面の「通知設定」または「プロフィール設定」をクリック"),
        ("「LINE連携コードを発行」ボタンをクリック",
         "画面に8桁の連携コードが表示されます"),
        ("LINE アプリでコードを送信",
         "LINEで「シフト管理 公式アカウント」に友達追加し、連携コードを送信します"),
        ("連携完了を確認",
         "「LINE連携済み」と表示されれば設定完了です。以後LINEで通知を受け取れます"),
    ])

    # 17. 通知設定
    add_two_column_slide(prs, "通知設定",
        "LINE通知", [
            "連携済みの場合、以下のタイミングで通知:",
            "・ 募集開始時",
            "・ 締切リマインダー（締切3日前）",
            "・ 承認 / 返却 時",
            "・ 確定シフト公開時",
            "設定画面でON/OFFを切り替え可能",
        ],
        "メール通知", [
            "登録メールアドレスに通知を送信:",
            "・ LINE通知と同じタイミング",
            "・ メールアドレスはプロフィール設定で変更可",
            "・ 設定画面でON/OFFを切り替え可能",
            "LINE・メール両方ONにすることも可能",
        ])

    # 18. ルールと規則の概要
    add_content_slide(prs, "ルールと規則の概要", [
        "提出単位と期間:",
        (1, "シフト希望は【2週間単位】で提出します"),
        (1, "期間の開始日・終了日・締切は管理者が設定します"),
        "締切の厳守:",
        (1, "提出締切を必ず守ってください"),
        (1, "締切後の変更は管理者の許可が必要です"),
        "ヒートマップ赤表示の日の協力義務:",
        (1, "赤表示の日は人員が危機的に不足しています"),
        (1, "可能な限り「出勤」への変更にご協力ください"),
        "その他:",
        (1, "代理休みは必ず代替者と事前に合意してから入力してください"),
        (1, "不明点は管理者に直接お問い合わせください"),
    ])

    # 19. よくある質問 (Q&A)
    add_content_slide(prs, "よくある質問 (Q&A)", [
        "Q: 締切を過ぎてしまいました。どうすればよいですか？",
        (1, "A: 管理者に直接連絡してください。管理者が代理入力または期間を再オープンします"),
        "Q: 「返却」されました。何をすればよいですか？",
        (1, "A: ステータスをクリックして管理者メモを確認し、内容を修正して再提出してください"),
        "Q: パスワードを忘れました",
        (1, "A: 管理者に連絡し、パスワードリセットを依頼してください"),
        "Q: スマホで表示が崩れます",
        (1, "A: Chrome ブラウザの最新版をお試しください。または PC でのアクセスをお試しください"),
        "Q: LINE通知が届きません",
        (1, "A: LINE連携の設定を確認してください。未連携の場合は手順16を参照してください"),
    ])

    prs.save(output_path)
    print(f"✓ {output_path} を保存しました ({len(prs.slides)} スライド)")


# ---------------------------------------------------------------------------
# 管理者向けマニュアル
# ---------------------------------------------------------------------------

def build_admin_manual(output_path="admin_manual.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    C = RGBColor(0x06, 0x4E, 0x3B)  # 管理者向け: 深緑

    # 1. タイトル
    add_title_slide(prs,
                    "管理者向け\nシフト管理システム 操作マニュアル",
                    "～ 期間管理・承認・スケジュール確定の手引き ～",
                    color=C)

    # 2. 目次
    toc = [
        "システム概要と管理者権限",
        "期間の作成と管理",
        "従業員への通知送信",
        "提出一覧の確認",
        "提出内容の承認・返却",
        "ダッシュボード",
        "日別シフト割り当て",
        "代理入力",
        "従業員プロファイル管理",
        "会社休日の設定",
        "CSV 出力",
        "ステータス遷移まとめ",
    ]
    add_toc_slide(prs, toc)

    # 3. システム概要と管理者権限
    add_two_column_slide(prs, "システム概要と管理者権限",
        "管理者ができること", [
            "シフト希望期間の作成・管理",
            "従業員への通知送信",
            "提出内容の承認 / 返却",
            "シフトの確定と公開",
            "日別シフト割り当て",
            "代理入力（締切後も可）",
            "従業員プロファイルの編集",
            "会社休日の設定",
            "CSV データ出力",
        ],
        "従業員との違い", [
            "全従業員の提出状況を閲覧可能",
            "期間のステータスを変更可能",
            "提出締切後も入力・修正が可能",
            "管理者メモを記入して返却可能",
            "ドライバー設定など属性を変更可能",
            "システム全体の設定を変更可能",
        ],
        color=C)

    # 4. 期間の作成と管理
    add_steps_slide(prs, "期間の作成と管理", [
        ("管理者メニュー → 「期間管理」を開く",
         "サイドバーまたはヘッダーの管理者メニューからアクセスします"),
        ("「新規期間作成」ボタンをクリック",
         "開始日・終了日・提出締切日時を入力します（通常2週間単位）"),
        ("期間を保存し OPEN ステータスに設定",
         "保存後、ステータスを OPEN にすると従業員への募集が開始されます"),
        ("締切後に REVIEW ステータスへ変更",
         "締切後、提出内容を確認しながら承認・返却を行います"),
        ("全員承認後に FIXED → PUBLISHED へ変更",
         "PUBLISHED にすると従業員が確定スケジュールを閲覧できます"),
    ], color=C)

    # 4b. ステータス遷移
    add_period_status_slide(prs, color=C)

    # 5. 従業員への通知送信
    add_content_slide(prs, "従業員への通知送信", [
        "通知の種類:",
        (1, "【募集開始通知】: 期間を OPEN にしたタイミングで一斉送信"),
        (1, "【リマインダー通知】: 締切3日前、または任意のタイミングで送信"),
        (1, "【個別通知】: 特定の従業員にのみ送信することも可能"),
        "送信手順:",
        (1, "対象期間の「通知」ボタンをクリック"),
        (1, "送信対象（全員 / 未提出者のみ / 個別）を選択"),
        (1, "メッセージを確認して「送信」ボタンをクリック"),
        "通知チャネル: LINE / メール（従業員の設定に応じて自動選択）",
    ], color=C)

    # 6. 提出一覧の確認
    add_content_slide(prs, "提出一覧の確認", [
        "「提出一覧」画面の開き方:",
        (1, "対象期間の詳細画面から「提出一覧」タブを選択"),
        "一覧で確認できる情報:",
        (1, "従業員名・提出日時・ステータス（提出済 / 未提出 / 返却済）"),
        (1, "提出内容の簡易サマリー（出勤希望日数 など）"),
        "フィルタ機能:",
        (1, "「未提出者のみ表示」で締切前のフォローアップが簡単"),
        (1, "勤務形態別（フルタイム / パート）でフィルタリング可能"),
        "各行をクリックすると提出内容の詳細が確認できます",
    ], color=C)

    # 7. 提出内容の承認・返却
    add_steps_slide(prs, "提出内容の承認・返却", [
        ("提出一覧から対象従業員の行をクリック",
         "提出された日程・理由・コメントが一覧表示されます"),
        ("内容を確認する",
         "ヒートマップで人員バランスを見ながら各日程を確認します"),
        ("「承認」または「返却」を選択",
         "問題なければ「承認」。修正が必要な場合は「返却」を選択します"),
        ("返却の場合: 管理者メモを入力",
         "修正してほしい内容を具体的に記載してください。従業員に通知が届きます"),
    ], color=C)

    # 8. ダッシュボード
    add_content_slide(prs, "ダッシュボード", [
        "ダッシュボードで確認できる情報:",
        (1, "スタッフィング状況: 日別の出勤予定者数と目標人数の比較"),
        (1, "危険日のハイライト: 目標を下回る日が赤でハイライト表示"),
        (1, "ドライバー不足日: ドライバー要員が不足している日がオレンジ表示"),
        (1, "未提出者数: まだ提出していない従業員数のカウント"),
        "活用方法:",
        (1, "赤ハイライトの日付を確認し、該当日の承認判断に活用"),
        (1, "未提出者が多い場合はリマインダー通知を送信"),
        (1, "特定日の人数が不足している場合は、休み申請を返却して調整"),
    ], color=C)

    # 9. 日別シフト割り当て
    add_steps_slide(prs, "日別シフト割り当て", [
        ("ダッシュボードまたはカレンダーから対象日をクリック",
         "その日の全売り場と割当状況が表示されます"),
        ("売り場（ポジション）を選択",
         "割り当てたい売り場の「＋」ボタンをクリックします"),
        ("従業員を選択して割り当て",
         "出勤希望者リストから選択。ドラッグ＆ドロップでも移動できます"),
        ("保存する",
         "割り当てを完了したら「保存」ボタンをクリックします"),
    ], color=C)

    # 10. 代理入力
    add_content_slide(prs, "代理入力（締切後の管理者による入力）", [
        "代理入力が必要なケース:",
        (1, "従業員が締切に間に合わなかった場合"),
        (1, "従業員がシステムを利用できない状況の場合"),
        "代理入力の手順:",
        (1, "提出一覧で対象従業員の行を確認"),
        (1, "「代理入力」ボタンをクリック（管理者のみ表示）"),
        (1, "通常の提出画面と同様に日程を入力"),
        (1, "「管理者代理入力」のラベルが付いて保存されます"),
        "注意: 代理入力後も従業員に通知が送信されます。入力内容を本人に伝えてください",
    ], color=C)

    # 11. 従業員プロファイル管理
    add_two_column_slide(prs, "従業員プロファイル管理",
        "編集できる項目", [
            "勤務形態（フルタイム / パート 等）",
            "固定曜日（毎週固定で出勤する曜日）",
            "ドライバー設定（ON / OFF）",
            "メールアドレス",
            "LINE連携状態の確認・リセット",
            "通知設定（LINE / メール）",
        ],
        "編集手順", [
            "管理者メニュー → 「従業員管理」",
            "対象従業員の行をクリック",
            "「編集」ボタンをクリック",
            "各項目を変更して保存",
            "変更内容は即時反映されます",
        ],
        color=C)

    # 12. 会社休日の設定
    add_steps_slide(prs, "会社休日の設定", [
        ("管理者メニュー → 「会社休日設定」を開く",
         "年間カレンダーが表示されます"),
        ("休日として追加したい日付をクリック",
         "クリックするとその日が会社休日としてマーク（グレー表示）されます"),
        ("休日の削除",
         "設定済みの休日をクリックすると解除できます"),
        ("保存する",
         "「保存」ボタンをクリックすると全従業員のシフト画面に反映されます"),
    ], color=C)

    # 13. CSV 出力
    add_steps_slide(prs, "CSV 出力", [
        ("対象期間の詳細画面を開く",
         "CSVで出力したい期間を選択します"),
        ("「CSV ダウンロード」ボタンをクリック",
         "画面上部または提出一覧タブ内にあります"),
        ("出力内容を選択（オプション）",
         "「全従業員」「提出済み」「承認済み」などフィルタを選べる場合があります"),
        ("ファイルをダウンロード",
         "ブラウザのダウンロードフォルダに CSV ファイルが保存されます"),
    ], color=C)

    prs.save(output_path)
    print(f"✓ {output_path} を保存しました ({len(prs.slides)} スライド)")


# ---------------------------------------------------------------------------
# エントリポイント
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    out_dir = os.path.dirname(os.path.abspath(__file__))
    # プロジェクトルートに出力
    root_dir = os.path.dirname(out_dir)

    emp_path = os.path.join(root_dir, "employee_manual.pptx")
    adm_path = os.path.join(root_dir, "admin_manual.pptx")

    print("シフト管理システム 操作マニュアルを生成します...")
    build_employee_manual(emp_path)
    build_admin_manual(adm_path)
    print("\n完了！ プロジェクトルートに以下のファイルが生成されました:")
    print(f"  {emp_path}")
    print(f"  {adm_path}")
